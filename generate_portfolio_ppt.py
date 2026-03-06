from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "Jupiter_Cyber_Defence_Client_Deck.pptx"

logo = ROOT / "jupiter_logo.jpeg"
branding = ROOT / "branding.jpeg"
gdpr = ROOT / "GDPR.png"
hipaa = ROOT / "HIPAA.png"
pdpl = ROOT / "PDPL.png"


COLOR_DARK = (7, 13, 26)
COLOR_MID = (12, 20, 40)
COLOR_LIGHT = (244, 247, 255)
COLOR_ACCENT = (122, 162, 255)
COLOR_ACCENT_2 = (158, 132, 255)
COLOR_TEXT_LIGHT = (238, 243, 255)
COLOR_TEXT_DARK = (25, 38, 67)
COLOR_TEXT_MUTED = (127, 148, 192)
COLOR_BORDER = (213, 222, 245)


def rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*color)


def add_bg(slide, width, height, color: tuple[int, int, int]) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()


def add_circle(slide, left, top, size, color, transparency=0.0):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = rgb(color)
    circle.fill.transparency = transparency
    circle.line.fill.background()


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=20,
    bold=False,
    color=COLOR_TEXT_DARK,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    p.alignment = align
    return box


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    items,
    size=16,
    color=COLOR_TEXT_DARK,
    bullet_color=COLOR_ACCENT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.font.bold = False
        p.space_after = Pt(7)
        p.space_before = Pt(0)
        p.bullet = True
        p.runs[0].font.color.rgb = rgb(color)
    box.fill.background()
    bullet_icon = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left - Inches(0.02), top - Inches(0.02), Inches(0.001), Inches(0.001))
    bullet_icon.fill.solid()
    bullet_icon.fill.fore_color.rgb = rgb(bullet_color)
    bullet_icon.line.fill.background()


def add_logo(slide, left, top, width=Inches(0.45), height=Inches(0.45)):
    if logo.exists():
        slide.shapes.add_picture(str(logo), left, top, width=width, height=height)


def style_cell(target, text, *, bold=False, size=12, color=COLOR_TEXT_DARK, align=PP_ALIGN.CENTER):
    tf = target.text_frame if hasattr(target, "text_frame") else target
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = rgb(color)
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)


def cover_slide(prs, width, height):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, width, height, COLOR_DARK)

    add_circle(slide, Inches(-1.2), Inches(-1.0), Inches(3.2), COLOR_ACCENT, 0.75)
    add_circle(slide, Inches(11.3), Inches(-0.8), Inches(2.8), COLOR_ACCENT_2, 0.8)

    add_logo(slide, Inches(0.6), Inches(0.45))
    add_textbox(slide, Inches(1.1), Inches(0.45), Inches(4.5), Inches(0.5), "Jupiter Cyber Defence", size=16, bold=True, color=COLOR_TEXT_LIGHT)

    add_textbox(
        slide,
        Inches(0.6),
        Inches(1.25),
        Inches(7.0),
        Inches(1.2),
        "Client Presentation Deck",
        size=40,
        bold=True,
        color=COLOR_TEXT_LIGHT,
    )

    add_textbox(
        slide,
        Inches(0.6),
        Inches(2.3),
        Inches(7.0),
        Inches(1.0),
        "From reactive checks to continuous cyber confidence",
        size=18,
        bold=False,
        color=(185, 198, 234),
    )

    bullets = [
        "Clear risk visibility and prioritized remediation",
        "Compliance mapping for GDPR, PDPL, and HIPAA",
        "Audit-ready reporting for leadership and partners",
    ]
    add_bullets(
        slide,
        Inches(0.8),
        Inches(3.2),
        Inches(6.1),
        Inches(2.1),
        bullets,
        size=14,
        color=(227, 235, 255),
    )

    if branding.exists():
        slide.shapes.add_picture(str(branding), Inches(7.25), Inches(1.2), width=Inches(5.4), height=Inches(4.8))

    cta = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.75))
    cta.fill.solid()
    cta.fill.fore_color.rgb = rgb(COLOR_MID)
    cta.line.color.rgb = rgb((92, 118, 173))
    style_cell(cta.text_frame, "Book your free security snapshot: jupitercyberdefence.com", bold=True, size=16, color=COLOR_TEXT_LIGHT)


def challenge_solution_slide(prs, width, height):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, width, height, COLOR_LIGHT)

    add_logo(slide, Inches(0.6), Inches(0.38))
    add_textbox(slide, Inches(1.1), Inches(0.35), Inches(8.0), Inches(0.5), "Jupiter Cyber Defence", size=14, bold=True)
    add_textbox(slide, Inches(0.6), Inches(0.95), Inches(11.8), Inches(0.8), "Why Clients Choose Jupiter", size=30, bold=True)

    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(1.9), Inches(5.95), Inches(4.9))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = rgb((255, 255, 255))
    left_card.line.color.rgb = rgb(COLOR_BORDER)

    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.75), Inches(1.9), Inches(5.95), Inches(4.9))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = rgb((255, 255, 255))
    right_card.line.color.rgb = rgb(COLOR_BORDER)

    add_textbox(slide, Inches(1.0), Inches(2.2), Inches(5.2), Inches(0.6), "Common Security Pain Points", size=18, bold=True)
    add_bullets(
        slide,
        Inches(1.0),
        Inches(2.8),
        Inches(5.2),
        Inches(3.6),
        [
            "Too many tools, unclear priorities",
            "Limited visibility on external exposure",
            "Compliance pressure without clear mapping",
            "Reports that are too technical for leadership",
        ],
        size=14,
    )

    add_textbox(slide, Inches(7.1), Inches(2.2), Inches(5.2), Inches(0.6), "What Jupiter Delivers", size=18, bold=True)
    add_bullets(
        slide,
        Inches(7.1),
        Inches(2.8),
        Inches(5.2),
        Inches(3.6),
        [
            "One clear view of risk posture",
            "Business-prioritized remediation roadmap",
            "Regulatory mapping with gap visibility",
            "Executive-ready audit evidence and reporting",
        ],
        size=14,
    )


def services_slide(prs, width, height):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, width, height, COLOR_DARK)

    add_logo(slide, Inches(0.6), Inches(0.38))
    add_textbox(slide, Inches(1.1), Inches(0.35), Inches(8.0), Inches(0.5), "Jupiter Cyber Defence", size=14, bold=True, color=COLOR_TEXT_LIGHT)
    add_textbox(slide, Inches(0.6), Inches(0.95), Inches(11.8), Inches(0.8), "Core Services", size=30, bold=True, color=COLOR_TEXT_LIGHT)

    services_left = [
        "Cyber Security Posture Assessment",
        "External Attack Surface Assessment",
        "Continuous Risk Monitoring",
        "Cloud Security Readiness (AWS/Azure/GCP)",
        "Identity & Access Security Review",
        "Backup & Ransomware Readiness",
    ]
    services_right = [
        "Compliance Mapping (GDPR, PDPL, HIPAA)",
        "Audit-Ready Documentation",
        "Executive Security Reporting",
        "Remediation Prioritization",
        "Incident Readiness Review",
        "Third-Party Risk Review",
    ]

    panel_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(1.95), Inches(5.95), Inches(4.95))
    panel_left.fill.solid()
    panel_left.fill.fore_color.rgb = rgb((15, 26, 49))
    panel_left.line.color.rgb = rgb((81, 102, 153))

    panel_right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.75), Inches(1.95), Inches(5.95), Inches(4.95))
    panel_right.fill.solid()
    panel_right.fill.fore_color.rgb = rgb((15, 26, 49))
    panel_right.line.color.rgb = rgb((81, 102, 153))

    add_bullets(slide, Inches(1.0), Inches(2.25), Inches(5.45), Inches(4.4), services_left, size=13, color=COLOR_TEXT_LIGHT, bullet_color=COLOR_ACCENT)
    add_bullets(slide, Inches(7.1), Inches(2.25), Inches(5.45), Inches(4.4), services_right, size=13, color=COLOR_TEXT_LIGHT, bullet_color=COLOR_ACCENT)


def compliance_slide(prs, width, height):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, width, height, COLOR_LIGHT)

    add_logo(slide, Inches(0.6), Inches(0.38))
    add_textbox(slide, Inches(1.1), Inches(0.35), Inches(8.0), Inches(0.5), "Jupiter Cyber Defence", size=14, bold=True)
    add_textbox(slide, Inches(0.6), Inches(0.95), Inches(11.8), Inches(0.8), "Compliance and Evidence", size=30, bold=True)

    add_textbox(
        slide,
        Inches(0.6),
        Inches(1.6),
        Inches(11.9),
        Inches(0.6),
        "Map controls against GDPR, PDPL, and HIPAA requirements with prioritized actions.",
        size=16,
        color=COLOR_TEXT_MUTED,
    )

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(2.2), Inches(12.0), Inches(4.1))
    card.fill.solid()
    card.fill.fore_color.rgb = rgb((255, 255, 255))
    card.line.color.rgb = rgb(COLOR_BORDER)

    logo_paths = [gdpr, hipaa, pdpl]
    x = 1.35
    for p in logo_paths:
        frame = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.55), Inches(1.2), Inches(1.2))
        frame.fill.solid()
        frame.fill.fore_color.rgb = rgb((245, 248, 255))
        frame.line.color.rgb = rgb(COLOR_BORDER)
        if p.exists():
            slide.shapes.add_picture(str(p), Inches(x + 0.18), Inches(2.73), width=Inches(0.84), height=Inches(0.84))
        x += 1.55

    add_bullets(
        slide,
        Inches(0.95),
        Inches(4.1),
        Inches(5.4),
        Inches(1.9),
        [
            "Control mapping and gap status",
            "Priority remediation by impact and effort",
            "Traceable ownership and progress view",
        ],
        size=13,
    )

    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.65), Inches(2.55), Inches(5.55), Inches(3.45))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = rgb((238, 243, 255))
    right_box.line.color.rgb = rgb((205, 218, 245))
    add_textbox(slide, Inches(6.95), Inches(2.85), Inches(4.95), Inches(0.5), "Audit-Ready Outputs", size=17, bold=True)
    add_bullets(
        slide,
        Inches(6.95),
        Inches(3.35),
        Inches(4.95),
        Inches(2.4),
        [
            "Executive-friendly reporting",
            "Heatmaps and action ownership",
            "Evidence pack for audits and partners",
        ],
        size=13,
    )


def pricing_slide(prs, width, height):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, width, height, COLOR_DARK)

    add_logo(slide, Inches(0.6), Inches(0.38))
    add_textbox(slide, Inches(1.1), Inches(0.35), Inches(8.0), Inches(0.5), "Jupiter Cyber Defence", size=14, bold=True, color=COLOR_TEXT_LIGHT)
    add_textbox(slide, Inches(0.6), Inches(0.95), Inches(11.8), Inches(0.8), "Package Pricing Matrix", size=30, bold=True, color=COLOR_TEXT_LIGHT)
    add_textbox(slide, Inches(0.6), Inches(1.6), Inches(11.9), Inches(0.6), "Indicative package ranges by country", size=14, color=(189, 202, 237))

    table_shape = slide.shapes.add_table(4, 5, Inches(0.65), Inches(2.2), Inches(12.0), Inches(3.2))
    table = table_shape.table

    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.5)
    table.columns[3].width = Inches(2.5)
    table.columns[4].width = Inches(2.5)

    headers = ["Package", "Oman", "UAE", "KSA", "EU"]
    rows = [
        ["Package 1", "OMR 80–100", "AED 300–500", "SAR 300–500", "EUR 100–150"],
        ["Package 2", "OMR 200–400", "AED 2,000–4,000", "SAR 2,000–5,000", "EUR 500–800"],
        ["Package 3", "OMR 500–1,000", "AED 4,500–7,500", "SAR 6,000–12,000", "EUR 1,000–1,500"],
    ]

    for col, text in enumerate(headers):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb((33, 52, 90))
        style_cell(cell, text, bold=True, size=13, color=COLOR_TEXT_LIGHT)

    for r_index, row_vals in enumerate(rows, start=1):
        for c_index, text in enumerate(row_vals):
            cell = table.cell(r_index, c_index)
            cell.fill.solid()
            if r_index % 2 == 0:
                cell.fill.fore_color.rgb = rgb((247, 250, 255))
                txt_color = (38, 56, 95)
            else:
                cell.fill.fore_color.rgb = rgb((232, 240, 255))
                txt_color = (27, 44, 77)
            style_cell(cell, text, bold=(c_index == 0), size=12, color=txt_color)

    add_textbox(
        slide,
        Inches(0.65),
        Inches(5.7),
        Inches(12.0),
        Inches(0.7),
        "Packages can be tailored by asset volume, technical scope, and required reporting depth.",
        size=12,
        color=(178, 194, 231),
        align=PP_ALIGN.CENTER,
    )


def process_outcomes_slide(prs, width, height):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, width, height, COLOR_LIGHT)

    add_logo(slide, Inches(0.6), Inches(0.38))
    add_textbox(slide, Inches(1.1), Inches(0.35), Inches(8.0), Inches(0.5), "Jupiter Cyber Defence", size=14, bold=True)
    add_textbox(slide, Inches(0.6), Inches(0.95), Inches(11.8), Inches(0.8), "How Engagement Works", size=30, bold=True)

    steps = [
        ("1", "Understand", "Map domains, cloud, apps, and data flows."),
        ("2", "Assess", "Review posture, exposure, identity, and backups."),
        ("3", "Prioritize", "Deliver a 30–60 day remediation roadmap."),
    ]

    x = 0.8
    for num, title, desc in steps:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(3.95), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = rgb((255, 255, 255))
        card.line.color.rgb = rgb(COLOR_BORDER)

        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.15), Inches(2.15), Inches(0.58), Inches(0.58))
        badge.fill.solid()
        badge.fill.fore_color.rgb = rgb(COLOR_ACCENT)
        badge.line.fill.background()
        style_cell(badge.text_frame, num, bold=True, size=16, color=(255, 255, 255))

        add_textbox(slide, Inches(x + 0.85), Inches(2.2), Inches(2.9), Inches(0.5), title, size=17, bold=True)
        add_textbox(slide, Inches(x + 0.3), Inches(2.85), Inches(3.45), Inches(1.1), desc, size=13, color=(58, 81, 126))
        x += 4.15

    outcome_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.6), Inches(12.0), Inches(2.3))
    outcome_box.fill.solid()
    outcome_box.fill.fore_color.rgb = rgb((14, 23, 44))
    outcome_box.line.color.rgb = rgb((76, 96, 145))

    add_textbox(slide, Inches(1.15), Inches(4.9), Inches(4.0), Inches(0.5), "Business Outcomes", size=18, bold=True, color=COLOR_TEXT_LIGHT)
    add_bullets(
        slide,
        Inches(1.15),
        Inches(5.35),
        Inches(5.4),
        Inches(1.35),
        [
            "Faster remediation cycles",
            "Stronger customer and partner trust",
            "Clear weekly risk trend visibility",
        ],
        size=13,
        color=COLOR_TEXT_LIGHT,
        bullet_color=COLOR_ACCENT,
    )

    add_textbox(slide, Inches(7.2), Inches(4.9), Inches(5.2), Inches(0.5), "Industries", size=18, bold=True, color=COLOR_TEXT_LIGHT)
    chips = ["Healthcare", "Logistics", "SaaS", "Education", "Professional Services"]
    cx = 7.2
    cy = 5.35
    for chip in chips:
        w = max(1.35, len(chip) * 0.07)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(cy), Inches(w), Inches(0.42))
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb((24, 36, 68))
        shp.line.color.rgb = rgb((93, 116, 171))
        style_cell(shp.text_frame, chip, bold=False, size=11, color=(218, 229, 255))
        cx += w + 0.18
        if cx > 11.8:
            cx = 7.2
            cy += 0.52


def closing_slide(prs, width, height):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, width, height, COLOR_DARK)

    add_circle(slide, Inches(-1.1), Inches(4.8), Inches(2.9), COLOR_ACCENT, 0.75)
    add_circle(slide, Inches(11.2), Inches(-0.6), Inches(2.5), COLOR_ACCENT_2, 0.8)

    add_logo(slide, Inches(0.6), Inches(0.42), width=Inches(0.55), height=Inches(0.55))
    add_textbox(slide, Inches(1.25), Inches(0.42), Inches(6.0), Inches(0.6), "Jupiter Cyber Defence", size=16, bold=True, color=COLOR_TEXT_LIGHT)

    add_textbox(slide, Inches(0.6), Inches(1.45), Inches(12.0), Inches(0.9), "Book Your Free Security Snapshot", size=40, bold=True, color=COLOR_TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.6), Inches(2.35), Inches(12.0), Inches(0.7), "Turn security complexity into business clarity.", size=18, color=(189, 202, 237), align=PP_ALIGN.CENTER)

    if branding.exists():
        slide.shapes.add_picture(str(branding), Inches(4.5), Inches(3.1), width=Inches(4.3), height=Inches(2.2))

    contact = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.15), Inches(5.65), Inches(9.0), Inches(1.05))
    contact.fill.solid()
    contact.fill.fore_color.rgb = rgb(COLOR_MID)
    contact.line.color.rgb = rgb((97, 120, 176))
    style_cell(contact.text_frame, "info@jupitercyberdefence.com  |  jupitercyberdefence.com", bold=True, size=18, color=COLOR_TEXT_LIGHT)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    width = prs.slide_width
    height = prs.slide_height

    cover_slide(prs, width, height)
    challenge_solution_slide(prs, width, height)
    services_slide(prs, width, height)
    compliance_slide(prs, width, height)
    pricing_slide(prs, width, height)
    process_outcomes_slide(prs, width, height)
    closing_slide(prs, width, height)

    prs.save(str(OUTPUT))
    print(f"Created: {OUTPUT.resolve()}")


if __name__ == "__main__":
    main()
